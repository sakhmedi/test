from __future__ import annotations

import io

CHUNK_SIZE = 800  # chars
OVERLAP = 100     # chars


def extract_pages(filename: str, data: bytes) -> list[dict]:
    """Return list of {"page": int, "text": str} from file bytes."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in ("jpg", "jpeg", "png", "tiff"):
        # Raw image bytes can't be meaningfully decoded as text.
        # Callers should use VisionClient first to convert to a .txt file.
        return []

    elif ext == "pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            pages = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append({"page": i, "text": text})
            return pages
        except Exception:
            return []

    elif ext in ("docx", "doc"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(data))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return [{"page": 1, "text": text}] if text.strip() else []
        except Exception:
            return []

    elif ext in ("xlsx", "xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join("" if v is None else str(v) for v in row)
                    if line.strip():
                        parts.append(line)
            text = "\n".join(parts)
            return [{"page": 1, "text": text}] if text.strip() else []
        except Exception:
            return []

    elif ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            pages = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs)
                            if line.strip():
                                texts.append(line)
                if texts:
                    pages.append({"page": i, "text": "\n".join(texts)})
            return pages
        except Exception:
            return []

    else:
        # Plain text (covers .txt, _ocr.txt, and any other text files)
        return [{"page": 1, "text": data.decode("utf-8", errors="replace")}]


def chunk_pages(pages: list[dict]) -> list[dict]:
    """Return list of {"page": int, "text": str} chunks."""
    chunks = []
    for p in pages:
        text = p["text"].strip()
        start = 0
        while start < len(text):
            end = start + CHUNK_SIZE
            chunks.append({"page": p["page"], "text": text[start:end]})
            start += CHUNK_SIZE - OVERLAP
    return [c for c in chunks if len(c["text"].strip()) >= 30]
